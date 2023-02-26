import pptx
from googletrans import translator
import cv2
import streamlit as st

ppt_file = pptx.Presentation("presentation.pptx")
translator = Translator(service_urls = ["translate.google.com"])

def new_f():
  for slide in ppt_file.slides:
    for shape in slide.shapes:
      if shape.has_text_frame:
        text = shape.text_frame.text
        translated_text = translator.translate(text,dest='hi').text
        shape.text_frame.text = translated_text
  ppt_file.save("translated_presentation.pptx")
  
st.title("now we'll translate")
checking = st.button("Run)

if checking == True:
  new_f()


